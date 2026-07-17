from __future__ import annotations

import sys
from pathlib import Path


def safe_print(text: str) -> None:
    sys.stdout.buffer.write((text + "\n").encode("utf-8", errors="replace"))


def main() -> None:
    target = Path(sys.argv[1])
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover
        print(f"IMPORT_ERROR: {exc}")
        return

    prs = Presentation(str(target))
    safe_print(f"slides={len(prs.slides)}")
    for idx, slide in enumerate(prs.slides, start=1):
        safe_print(f"-- slide {idx} --")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip().replace("\n", " | ")
                if text:
                    safe_print(text[:500])


if __name__ == "__main__":
    main()
