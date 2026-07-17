from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation


UPDATES = {
    5: {
        4: "Train baseline tren bo 5 class de lay moc so sanh ban dau truoc khi toi uu bai toan.",
        6: "Rut gon ve 3 class: helmet, no_helmet, rider de tap trung dung muc tieu san pham.",
        8: "Fine-tune YOLOv8 voi train 2 stage, dieu chinh imgsz, epoch, patience, augmentation va kiem tra them TTA.",
        11: (
            "• Cau hinh train chinh da thu: thay doi imgsz 640 -> 768 -> 896, fine-tune tu checkpoint tot nhat, "
            "oversample no_helmet va danh gia bang standard val + TTA\n"
            "• Nhom cung thu YOLO11 va detector + classifier, nhung ket qua van thap hon YOLOv8 toi uu"
        ),
    },
    6: {
        5: "Train 5 class, danh gia de lay baseline cho Precision, Recall, mAP50 va mAP50-95",
        8: "Remap 5 class -> 3 class, fine-tune lai de giam nhieu va tang do phu hop bai toan",
        11: "Dieu chinh imgsz, epoch, patience, batch, augmentation, TTA va oversample no_helmet",
        14: "Thu YOLO11, hard-mining va 2-stage detector + classifier de so sanh",
    },
    7: {
        4: (
            "• Baseline 5 class: mAP50 = 65.1%, mAP50-95 = 30.2%\n"
            "• YOLOv8 fine-tune 3 class: Precision = 71.0%, Recall = 78.4%, mAP50 = 77.4%, mAP50-95 = 36.8%\n"
            "• Khi dung TTA: mAP50 = 77.9%, mAP50-95 = 36.8%\n"
            "• Muc cai thien so voi baseline: tang khoang 12.3 diem mAP50"
        ),
        5: "Kết luận: YOLOv8 fine-tune 3 class với train 2 stage là lựa chọn chính để demo và báo cáo.",
    },
}


def main() -> None:
    input_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2])

    prs = Presentation(str(input_path))

    for slide_number, shape_updates in UPDATES.items():
        slide = prs.slides[slide_number - 1]
        for shape_idx, new_text in shape_updates.items():
            slide.shapes[shape_idx].text = new_text

    # Clean duplicated conclusion line on slide 7 if present in this theme/layout.
    slide7 = prs.slides[6]
    for shape in slide7.shapes:
        if hasattr(shape, "text") and shape.text.strip() == "Kết luận: YOLOv8 fine-tune 3 class là lựa chọn chính để demo và báo cáo.":
            shape.text = ""

    prs.save(str(output_path))
    print(f"Wrote {output_path}")


if __name__ == "__main__":
    main()
