from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation


def safe_print(text: str) -> None:
    sys.stdout.buffer.write((text + "\n").encode("utf-8", errors="replace"))


def main() -> None:
    target = Path(sys.argv[1])
    prs = Presentation(str(target))
    for slide_idx, slide in enumerate(prs.slides, start=1):
        safe_print(f"== slide {slide_idx} ==")
        for idx, shape in enumerate(slide.shapes):
            name = getattr(shape, "name", "")
            text = getattr(shape, "text", "").strip().replace("\n", " | ")
            safe_print(f"[{idx}] {name}: {text[:400]}")


if __name__ == "__main__":
    main()
